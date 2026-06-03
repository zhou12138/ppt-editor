"""
PPTX 自然语言编辑器 MVP
用法: python3 pptx_editor.py <pptx文件> "自然语言指令"

无需 LLM API，纯规则引擎解析意图。
"""

import sys
import re
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ========== 颜色映射 ==========
COLOR_MAP = {
    "红": "FF0000", "红色": "FF0000",
    "蓝": "0000FF", "蓝色": "0000FF",
    "绿": "00FF00", "绿色": "00AA00",
    "黄": "FFFF00", "黄色": "FFD700",
    "黑": "000000", "黑色": "000000",
    "白": "FFFFFF", "白色": "FFFFFF",
    "灰": "888888", "灰色": "888888",
    "橙": "FF8C00", "橙色": "FF8C00",
    "紫": "800080", "紫色": "800080",
    "粉": "FF69B4", "粉色": "FF69B4",
    "深蓝": "003366", "浅蓝": "66CCFF",
    "深红": "8B0000", "深绿": "006400",
    "red": "FF0000", "blue": "0000FF", "green": "00AA00",
    "black": "000000", "white": "FFFFFF", "yellow": "FFD700",
    "orange": "FF8C00", "purple": "800080", "pink": "FF69B4",
    "gray": "888888", "grey": "888888",
}

# ========== PPTX 解析 ==========
def parse_pptx(path: str) -> dict:
    """解析 PPTX 为结构化描述"""
    prs = Presentation(path)
    result = {"slides": [], "slide_width": prs.slide_width, "slide_height": prs.slide_height}
    
    for si, slide in enumerate(prs.slides):
        slide_desc = {"index": si + 1, "layout": slide.slide_layout.name, "elements": []}
        
        for shape in slide.shapes:
            elem = {
                "id": shape.shape_id,
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left, "top": shape.top,
                "width": shape.width, "height": shape.height,
                "is_placeholder": shape.is_placeholder,
                "text": "",
            }
            
            if shape.is_placeholder:
                phf = shape.placeholder_format
                elem["ph_idx"] = phf.idx
                elem["ph_type"] = str(phf.type)
            
            if shape.has_text_frame:
                elem["text"] = shape.text_frame.text
                elem["paragraphs"] = []
                for para in shape.text_frame.paragraphs:
                    p = {"text": para.text, "runs": []}
                    for run in para.runs:
                        r = {
                            "text": run.text,
                            "font_name": run.font.name,
                            "font_size": round(run.font.size / 12700, 1) if run.font.size else None,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "color": str(run.font.color.rgb) if run.font.color and run.font.color.type else None,
                        }
                        p["runs"].append(r)
                    elem["paragraphs"].append(p)
            
            if shape.has_table:
                tbl = shape.table
                elem["table"] = [[cell.text for cell in row.cells] for row in tbl.rows]
            
            # 位置语义
            sw, sh = prs.slide_width, prs.slide_height
            cx = shape.left + shape.width / 2
            cy = shape.top + shape.height / 2
            h = "左" if cx < sw * 0.33 else ("右" if cx > sw * 0.67 else "中")
            v = "上" if cy < sh * 0.33 else ("下" if cy > sh * 0.67 else "中")
            elem["position_label"] = h + v
            
            slide_desc["elements"].append(elem)
        
        result["slides"].append(slide_desc)
    
    return result


def print_structure(desc: dict):
    """打印 PPTX 结构"""
    for slide in desc["slides"]:
        print(f"\n{'='*50}")
        print(f"📄 第 {slide['index']} 页 (布局: {slide['layout']})")
        print(f"{'='*50}")
        for e in slide["elements"]:
            ph = f" [占位符:{e.get('ph_type','')}]" if e["is_placeholder"] else ""
            pos = e["position_label"]
            text_preview = e["text"][:40].replace("\n", "↵") if e["text"] else "(无文本)"
            print(f"  [{e['id']}] {e['name']}{ph} ({pos})")
            print(f"       文本: {text_preview}")
            if "table" in e:
                print(f"       表格: {len(e['table'])}行×{len(e['table'][0])}列")


# ========== 意图解析（纯规则，不用 LLM） ==========
def parse_intent(instruction: str, pptx_desc: dict) -> list:
    """
    解析自然语言指令为操作列表
    返回: [{"action": ..., "slide": ..., "target": ..., "params": ...}, ...]
    """
    intents = []
    
    # ---- 提取页码 ----
    slide_num = None
    m = re.search(r'第(\d+)页', instruction)
    if m:
        slide_num = int(m.group(1))
    m = re.search(r'第(\d+)张', instruction)
    if m:
        slide_num = int(m.group(1))
    
    # ---- 提取目标 ----
    target = {}
    
    # 副标题（必须在标题之前匹配，因为"副标题"包含"标题"）
    if any(w in instruction for w in ["副标题", "subtitle"]):
        target["type"] = "subtitle"
    # 标题
    elif any(w in instruction for w in ["标题", "title", "题目"]):
        target["type"] = "title"
    # 正文/内容
    elif any(w in instruction for w in ["正文", "内容", "body", "文本"]):
        target["type"] = "body"
    # 表格
    elif "表格" in instruction or "table" in instruction.lower():
        target["type"] = "table"
    # 图片
    elif any(w in instruction for w in ["图片", "图", "picture", "image", "logo"]):
        target["type"] = "picture"
    
    # 位置描述
    for pos in ["左上", "右上", "左下", "右下", "居中", "中间"]:
        if pos in instruction:
            target["position"] = pos
    
    # 文本内容匹配（仅用于定位，排除"改成X"的X）
    for qm in re.finditer(r'[""「](.+?)[""」]', instruction):
        before = instruction[:qm.start()]
        if not re.search(r'(?:改|换|替换|变)[成为]?\s*$', before):
            target["text_match"] = qm.group(1)
            break
    
    # ---- 提取操作 + 参数 ----
    
    # 修改文本内容
    m = re.search(r'改[成为]?\s*[""「](.+?)[""」]', instruction)
    if not m:
        m = re.search(r'替换[成为]?\s*[""「](.+?)[""」]', instruction)
    if not m:
        m = re.search(r'换[成为]\s*[""「](.+?)[""」]', instruction)
    if m:
        intents.append({
            "action": "modify_text",
            "slide": slide_num,
            "target": target,
            "params": {"new_text": m.group(1)}
        })
    
    # 修改字号
    m = re.search(r'字号?\s*(?:改[成为]?|调[成为]?|设为)?\s*(\d+)', instruction)
    if m:
        intents.append({
            "action": "modify_font",
            "slide": slide_num,
            "target": target,
            "params": {"font_size": int(m.group(1))}
        })
    elif "大一点" in instruction or "放大" in instruction:
        intents.append({
            "action": "modify_font",
            "slide": slide_num,
            "target": target,
            "params": {"font_size_factor": 1.3}
        })
    elif "小一点" in instruction or "缩小" in instruction:
        intents.append({
            "action": "modify_font",
            "slide": slide_num,
            "target": target,
            "params": {"font_size_factor": 0.75}
        })
    
    # 加粗
    if any(w in instruction for w in ["加粗", "粗体", "bold"]):
        intents.append({
            "action": "modify_font",
            "slide": slide_num,
            "target": target,
            "params": {"bold": True}
        })
    
    # 取消加粗
    if any(w in instruction for w in ["取消加粗", "不加粗", "细体"]):
        intents.append({
            "action": "modify_font",
            "slide": slide_num,
            "target": target,
            "params": {"bold": False}
        })
    
    # 斜体
    if any(w in instruction for w in ["斜体", "italic"]):
        intents.append({
            "action": "modify_font",
            "slide": slide_num,
            "target": target,
            "params": {"italic": True}
        })
    
    # 颜色
    for color_name, color_hex in COLOR_MAP.items():
        if color_name in instruction:
            # 判断是"改成X色"还是"X色的元素"
            # 如果颜色词在"改成/换成"后面，是操作参数
            if re.search(rf'(?:改|换|变|调)[成为]?\s*{color_name}', instruction):
                intents.append({
                    "action": "modify_font",
                    "slide": slide_num,
                    "target": target,
                    "params": {"color": color_hex}
                })
            else:
                # 颜色词用于定位
                target["color"] = color_name
            break
    
    # 对齐
    align_map = {"居中": PP_ALIGN.CENTER, "左对齐": PP_ALIGN.LEFT, "右对齐": PP_ALIGN.RIGHT,
                 "两端对齐": PP_ALIGN.JUSTIFY}
    for align_name, align_val in align_map.items():
        if align_name in instruction:
            intents.append({
                "action": "modify_align",
                "slide": slide_num,
                "target": target,
                "params": {"alignment": align_val, "alignment_name": align_name}
            })
            break
    
    # 删除
    if any(w in instruction for w in ["删除", "删掉", "去掉", "移除"]):
        intents.append({
            "action": "delete",
            "slide": slide_num,
            "target": target,
            "params": {}
        })
    
    # 修改字体
    m = re.search(r'字体\s*(?:改[成为]?|换[成为]?|用)?\s*(\S+)', instruction)
    if m:
        font_name = m.group(1)
        if font_name not in ["大", "小", "一点"]:
            intents.append({
                "action": "modify_font",
                "slide": slide_num,
                "target": target,
                "params": {"font_name": font_name}
            })
    
    # 如果没解析出任何意图，返回"不理解"
    if not intents:
        intents.append({
            "action": "unknown",
            "raw": instruction,
            "target": target,
            "slide": slide_num,
        })
    
    return intents


# ========== 元素定位 ==========
def find_target(slide_desc: dict, target: dict) -> list:
    """在 slide 中查找匹配的元素"""
    candidates = slide_desc["elements"]
    
    # 按类型过滤
    if "type" in target:
        t = target["type"]
        if t == "title":
            candidates = [e for e in candidates if "TITLE" in e.get("ph_type", "") and "SUBTITLE" not in e.get("ph_type", "")]
        elif t == "subtitle":
            candidates = [e for e in candidates if "SUBTITLE" in e.get("ph_type", "") or e.get("ph_idx") == 1]
        elif t == "body":
            candidates = [e for e in candidates if "BODY" in e.get("ph_type", "")]
        elif t == "table":
            candidates = [e for e in candidates if "table" in e]
        elif t == "picture":
            candidates = [e for e in candidates if "PICTURE" in e.get("type", "")]
    
    # 按位置过滤
    if "position" in target:
        pos = target["position"]
        if pos == "中间" or pos == "居中":
            pos = "中中"
        candidates = [e for e in candidates if pos in e.get("position_label", "")]
    
    # 按文本内容匹配
    if "text_match" in target:
        query = target["text_match"]
        candidates = [e for e in candidates if query in e.get("text", "")]
    
    return candidates


# ========== 执行操作 ==========
def execute(pptx_path: str, intents: list, pptx_desc: dict, output_path: str):
    """执行操作"""
    prs = Presentation(pptx_path)
    changes = []
    
    for intent in intents:
        action = intent["action"]
        if action == "unknown":
            print(f"⚠️  无法理解: {intent['raw']}")
            continue
        
        # 确定 slide 范围
        if intent.get("slide"):
            slides = [(intent["slide"] - 1, prs.slides[intent["slide"] - 1])]
            slide_descs = [pptx_desc["slides"][intent["slide"] - 1]]
        else:
            slides = list(enumerate(prs.slides))
            slide_descs = pptx_desc["slides"]
        
        for (si, slide), slide_desc in zip(slides, slide_descs):
            targets = find_target(slide_desc, intent.get("target", {}))
            
            if not targets:
                print(f"⚠️  第{si+1}页: 没找到匹配的元素")
                continue
            
            if len(targets) > 1:
                print(f"ℹ️  第{si+1}页: 找到 {len(targets)} 个匹配元素:")
                for i, t in enumerate(targets):
                    text_preview = t["text"][:30] if t["text"] else "(无文本)"
                    print(f"     {i+1}. [{t['id']}] {t['name']} ({t['position_label']}) '{text_preview}'")
                print(f"     → 将对所有匹配元素执行操作")
            
            for target_elem in targets:
                shape = None
                for s in slide.shapes:
                    if s.shape_id == target_elem["id"]:
                        shape = s
                        break
                
                if not shape:
                    continue
                
                params = intent.get("params", {})
                
                if action == "modify_text":
                    if shape.has_text_frame:
                        old_text = shape.text_frame.text
                        # 保留格式，只改第一个 paragraph 第一个 run
                        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                            shape.text_frame.paragraphs[0].runs[0].text = params["new_text"]
                        else:
                            shape.text_frame.text = params["new_text"]
                        changes.append(f"✅ 第{si+1}页 [{target_elem['name']}] 文本: '{old_text[:20]}' → '{params['new_text'][:20]}'")
                
                elif action == "modify_font":
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if "font_size" in params:
                                    run.font.size = Pt(params["font_size"])
                                if "font_size_factor" in params:
                                    if run.font.size:
                                        run.font.size = int(run.font.size * params["font_size_factor"])
                                if "bold" in params:
                                    run.font.bold = params["bold"]
                                if "italic" in params:
                                    run.font.italic = params["italic"]
                                if "color" in params:
                                    run.font.color.rgb = RGBColor.from_string(params["color"])
                                if "font_name" in params:
                                    run.font.name = params["font_name"]
                        
                        desc_parts = []
                        if "font_size" in params: desc_parts.append(f"字号→{params['font_size']}pt")
                        if "font_size_factor" in params: desc_parts.append(f"字号×{params['font_size_factor']}")
                        if "bold" in params: desc_parts.append("加粗" if params["bold"] else "取消加粗")
                        if "italic" in params: desc_parts.append("斜体")
                        if "color" in params: desc_parts.append(f"颜色→#{params['color']}")
                        if "font_name" in params: desc_parts.append(f"字体→{params['font_name']}")
                        changes.append(f"✅ 第{si+1}页 [{target_elem['name']}] {', '.join(desc_parts)}")
                
                elif action == "modify_align":
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para.alignment = params["alignment"]
                        changes.append(f"✅ 第{si+1}页 [{target_elem['name']}] 对齐→{params['alignment_name']}")
                
                elif action == "delete":
                    sp = shape._element
                    sp.getparent().remove(sp)
                    changes.append(f"✅ 第{si+1}页 删除 [{target_elem['name']}]")
    
    if changes:
        prs.save(output_path)
        print(f"\n{'='*50}")
        print(f"📝 执行了 {len(changes)} 项修改:")
        for c in changes:
            print(f"   {c}")
        print(f"\n💾 已保存到: {output_path}")
    else:
        print("\n⚠️  没有执行任何修改")
    
    return changes


# ========== 主程序 ==========
def main():
    if len(sys.argv) < 2:
        print("PPTX 自然语言编辑器 MVP")
        print(f"用法: python3 {sys.argv[0]} <pptx文件> [指令]")
        print(f"      python3 {sys.argv[0]} <pptx文件> --inspect  (查看结构)")
        print()
        print("指令例子:")
        print('  "把第1页标题改成「Hello World」"')
        print('  "第2页标题字号36 加粗 红色"')
        print('  "删除第3页副标题"')
        print('  "标题改成蓝色"')
        sys.exit(0)
    
    pptx_path = sys.argv[1]
    if not Path(pptx_path).exists():
        print(f"❌ 文件不存在: {pptx_path}")
        sys.exit(1)
    
    # 解析结构
    desc = parse_pptx(pptx_path)
    
    # 查看模式
    if len(sys.argv) > 2 and sys.argv[2] == "--inspect":
        print_structure(desc)
        # 输出 JSON
        json_path = pptx_path.replace(".pptx", "_structure.json")
        # 清理不可序列化的值
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(desc, f, ensure_ascii=False, indent=2, default=str)
        print(f"\n📄 结构已导出: {json_path}")
        sys.exit(0)
    
    if len(sys.argv) < 3:
        print("请提供指令，或用 --inspect 查看结构")
        sys.exit(1)
    
    instruction = " ".join(sys.argv[2:])
    print(f"📝 指令: {instruction}")
    
    # 解析意图
    intents = parse_intent(instruction, desc)
    print(f"\n🔍 解析出 {len(intents)} 个意图:")
    for i, intent in enumerate(intents):
        print(f"   {i+1}. action={intent['action']}, target={intent.get('target', {})}")
        if intent.get("params"):
            print(f"      params={intent['params']}")
    
    # 执行
    output_path = pptx_path.replace(".pptx", "_modified.pptx")
    execute(pptx_path, intents, desc, output_path)


if __name__ == "__main__":
    main()
