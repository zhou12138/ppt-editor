"""Generate test PPTX - a professional starting template for demo_test to enhance."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# Set 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ====== Slide 1: Title page ======
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "PPT Editor Demo"
slide.placeholders[1].text = "COM Automation Test Suite"

# ====== Slide 2: Content with bullets ======
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Core Features"
body = slide.placeholders[1]
tf = body.text_frame
tf.text = "Text & Font Manipulation"
for item in ["Shape Styling & Effects", "Table Operations", "Chart Generation", "Animation & Transitions", "Export to PDF & Image"]:
    p = tf.add_paragraph()
    p.text = item

# ====== Slide 3: Table - product comparison ======
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Performance Metrics"
rows, cols = 5, 4
tbl_shape = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(2), Inches(10), Inches(4))
table = tbl_shape.table
headers = ["Metric", "Q1", "Q2", "Growth"]
data = [
    ["Revenue", "$1.2M", "$1.8M", "+50%"],
    ["Users", "45K", "82K", "+82%"],
    ["NPS Score", "68", "78", "+15%"],
    ["Uptime", "99.5%", "99.9%", "+0.4%"],
]
for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.text = h
    for run in cell.text_frame.paragraphs[0].runs:
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        table.cell(ri + 1, ci).text = val

# ====== Slide 4: Blank - for chart & freeform ======
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(1))
box.text_frame.text = "Data Visualization"
for run in box.text_frame.paragraphs[0].runs:
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2D, 0x3A, 0x4A)

# ====== Slide 5: Blank - visual effects showcase ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(1))
box.text_frame.text = "Visual Effects"
for run in box.text_frame.paragraphs[0].runs:
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2D, 0x3A, 0x4A)

# ====== Slide 6: Summary / Thank you ======
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Thank You"
slide.placeholders[1].text = "Powered by PPT Editor COM"

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "test_report.pptx")
prs.save(out)
print(f"[OK] test PPTX generated: {out} (6 slides, 16:9)")

# Generate test images
try:
    import struct, zlib
    def make_png(path, w, h, r, g, b):
        def chunk(ctype, data):
            c = ctype + data
            return struct.pack('>I', len(data)) + c + struct.pack('>I', zlib.crc32(c) & 0xffffffff)
        raw = b''
        for y in range(h):
            raw += b'\x00'
            for x in range(w):
                # gradient effect
                rr = min(255, r + int((x / w) * 60))
                gg = min(255, g + int((y / h) * 40))
                raw += bytes([rr, gg, b])
        return (b'\x89PNG\r\n\x1a\n' +
                chunk(b'IHDR', struct.pack('>IIBBBBB', w, h, 8, 2, 0, 0, 0)) +
                chunk(b'IDAT', zlib.compress(raw)) +
                chunk(b'IEND', b''))

    base = os.path.dirname(os.path.abspath(__file__))
    # Main test image - blue gradient
    with open(os.path.join(base, "test_img.png"), 'wb') as f:
        f.write(make_png("", 200, 150, 30, 100, 220))
    # Background image - dark gradient
    with open(os.path.join(base, "test_bg.png"), 'wb') as f:
        f.write(make_png("", 400, 300, 20, 30, 60))
    print("[OK] test images generated")
except Exception as e:
    print(f"[WARN] test image generation failed: {e}")
